# blueprints/create_pptx.py
import os
import traceback
import json
import io
from datetime import datetime
import zoneinfo

from azure.functions import Blueprint, HttpRequest, HttpResponse  # type: ignore
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Pt
from azure.storage.blob import BlobServiceClient, ContentSettings

# ──────────────────────────────────────────────────────────────
#  Blueprint
# ──────────────────────────────────────────────────────────────
create_pptx_bp = Blueprint()

# ──────────────────────────────────────────────────────────────
#  Azure OpenAI クライアント
# ──────────────────────────────────────────────────────────────
client = AzureOpenAI(
    api_version="2024-12-01-preview",
    azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
    api_key=os.environ["AZURE_OPENAI_KEY"],
)

SYSTEM_PROMPT = """
あなたはB2B向け資料作成のエキスパートです。
以下の JSON フォーマットで 5 枚分のスライド構成を日本語で出力してください。
[
  { "title": "タイトル1", "bullets": ["箇条書きA", "箇条書きB"] },
  ...
]
"""

# ──────────────────────────────────────────────────────────────
#  HTTP ルート
# ──────────────────────────────────────────────────────────────


@create_pptx_bp.route(route="auto_ppt", methods=["GET"])
def auto_ppt(req: HttpRequest) -> HttpResponse:
    """
    1) Azure OpenAI (gpt-4o) でスライド構成を生成
    2) python-pptx で .pptx を組み立て
    3) Blob Storage にアップロード
    4) ファイルをダウンロードさせつつアップロード結果をヘッダーで返す
    """

    USER_PROMPT = """ テーマ: 日本人夫婦が行くハワイ旅行 対象読者: 旅行会社の営業担当者 目的: ハワイの観光地を紹介する"""

    # ── 1. スライド構成を生成
    resp = client.chat.completions.create(
        model="gpt-4o",  # デプロイ名を合わせる
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": USER_PROMPT},
        ],
        max_completion_tokens=600,
    )
    slides = json.loads(resp.choices[0].message.content)

    # ── 2. PPTX を構築
    prs = Presentation()
    ts = datetime.now(zoneinfo.ZoneInfo("Asia/Tokyo")
                      ).strftime("%Y%m%d-%H%M%S")
    file_name = f"{ts}_auto_docs.pptx"

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = slides[0]["title"]
    cover.placeholders[1].text = f"Update Date – {ts}"

    for slide in slides[1:]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide["title"]
        tf = s.shapes.placeholders[1].text_frame
        tf.clear()
        for b in slide["bullets"]:
            p = tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    # ── 3. Blob Storage へアップロード
    try:
        blob_service = BlobServiceClient.from_connection_string(
            os.environ["BLOB_CONN"])
        container_name = "pptstorage"
        blob_path = f"generated/{file_name}"

        cc = blob_service.get_container_client(container_name)
        if not cc.exists():
            cc.create_container()

        cc.upload_blob(
            name=blob_path,
            data=buf.getvalue(),
            overwrite=True,
            content_settings=ContentSettings(
                content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
        )

        blob_url = (
            f"https://{blob_service.account_name}.blob.core.windows.net/"
            f"{container_name}/{blob_path}"
        )
        upload_status = f"Uploaded to Blob Storage: {blob_url}"
    except Exception as ex:
        upload_status = f"Blob upload failed: {ex}"

    # ── 4. HTTP 応答
    return HttpResponse(
        body=buf.getvalue(),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{file_name}"',
            "X-Upload-Status": upload_status,
        },
    )


@create_pptx_bp.route(route="ping", methods=["GET"])
def ping(req: HttpRequest) -> HttpResponse:
    """
    - OpenAI に 'ping' を送り、『pong』的な返答が返るかだけを見る
    - 失敗時は例外メッセージを 500 で返す
    """

    try:
        resp = client.chat.completions.create(
            model="gpt-4o",                # デプロイ名
            messages=[
                {"role": "system", "content": "You are a ping server. Reply with 'pong'."},
                {"role": "user",   "content": "ping"},
            ],
            max_completion_tokens=5,       # 極小
            timeout=15                     # 15 秒で十分
        )
        answer = resp.choices[0].message.content.strip()
        payload = {"status": "success", "answer": answer}
        return HttpResponse(json.dumps(payload, ensure_ascii=False),
                            mimetype="application/json", status_code=200)

    except Exception as e:
        err = {"status": "error", "detail": str(
            e), "trace": traceback.format_exc()}
        return HttpResponse(json.dumps(err, ensure_ascii=False),
                            mimetype="application/json", status_code=500)
